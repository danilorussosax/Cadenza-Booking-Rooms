#!/usr/bin/env python3
"""Genera la presentazione "Cadenza · Direzione" in formato `.pptx`,
apribile nativamente da Apple Keynote (Keynote → File → Apri →
selezionare il .pptx → File → Esporta come → Keynote).

Tutte le grafiche sono integrate:
- Slide background (paper o navy gradient simulato a strati)
- Card con angoli arrotondati e ombra opzionale
- Accent bar laterali, linee oro
- Icone lucide-react rasterizzate a 4x e incluse come PNG
- Screenshot reali della webapp (slide 5-10) come PNG embedded
- Text frame nativi con paragraph properties (font, size, color, align)

Output:
- /Users/danilorusso/Desktop/prenota-aule/Cadenza_Presentazione_Direzione.pptx
"""

from __future__ import annotations

import io
import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from lucide_render import render_icon

# ────────── paths ──────────
ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = Path(
    "/Users/danilorusso/Desktop/prenota-aule/Cadenza_Presentazione_Direzione.pptx"
)
ICON_PATH = next(
    p for p in [
        Path("/Users/danilorusso/Desktop/prenota-aule/conservatory-app/frontend/public/cadenza.png"),
        Path("/Users/danilorusso/Desktop/prenota-aule/cadenza.png"),
        Path("/Users/danilorusso/Desktop/prenota-aule/ALTRI FILE DI USO/cadenza.png"),
    ] if p.exists()
)

# ────────── pagina 1920×1080 punti EMU ──────────
# python-pptx usa EMU (English Metric Units): 914400 = 1 inch, 12700 = 1 pt
# Vogliamo slide 16:9 a 1920×1080 px.
SLIDE_W_PT = 1920
SLIDE_H_PT = 1080


def pt(x: float) -> int:
    """Punti → EMU."""
    return int(x * 12700)


# ────────── colour tokens (HSL→RGB dal webapp) ──────────
COLORS = {
    "NavyDeep":  RGBColor(15, 23, 42),
    "Navy":      RGBColor(28, 56, 115),
    "NavyLight": RGBColor(51, 87, 158),
    "Gold":      RGBColor(243, 148, 5),
    "GoldLight": RGBColor(250, 188, 80),
    "GoldPale":  RGBColor(252, 215, 155),
    "Paper":     RGBColor(248, 250, 252),
    "PaperDark": RGBColor(226, 232, 240),
    "Ink":       RGBColor(15, 23, 42),
    "Neutral":   RGBColor(100, 116, 139),
    "Dim":       RGBColor(148, 163, 184),
    "Green":     RGBColor(22, 163, 74),
    "GreenDark": RGBColor(16, 122, 56),
    "Red":       RGBColor(220, 38, 38),
    "RedDark":   RGBColor(175, 26, 26),
    "Violet":    RGBColor(139, 92, 246),
    "White":     RGBColor(255, 255, 255),
}


# ────────── helper drawing ──────────


def add_rect(slide, x, y, w, h, fill_name=None, *,
             stroke_name=None, stroke_w_pt=0, radius_pt=0):
    """Rettangolo (eventualmente con angoli arrotondati). Coordinate in pt."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_pt > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, pt(x), pt(y), pt(w), pt(h))
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # adjustment value: 0..0.5 → frazione del lato corto che diventa raggio
        adj_target = radius_pt / min(w, h) * 2
        shp.adjustments[0] = max(0, min(0.5, adj_target / 2))
    if fill_name is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[fill_name]
    if stroke_name is None or stroke_w_pt == 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = COLORS[stroke_name]
        shp.line.width = pt(stroke_w_pt)
    # remove default text inside auto-shape so it doesn't show placeholder
    shp.text_frame.word_wrap = True
    shp.text_frame.text = ""
    return shp


def add_oval(slide, x, y, w, h, fill_name=None, stroke_name=None,
             stroke_w_pt=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, pt(x), pt(y), pt(w), pt(h))
    if fill_name is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[fill_name]
    if stroke_name is None or stroke_w_pt == 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = COLORS[stroke_name]
        shp.line.width = pt(stroke_w_pt)
    shp.text_frame.text = ""
    return shp


def add_text(slide, x, y, w, h, text, *, size=22, color="Ink",
             bold=False, align="left", line_spacing=1.25):
    """TextBox con stile uniforme. `text` può contenere \\n per multi-line."""
    tb = slide.shapes.add_textbox(pt(x), pt(y), pt(w), pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    align_map = {"left": PP_ALIGN.LEFT,
                 "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT}
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS[color]
        run.font.name = "Helvetica"
    return tb


def add_image(slide, x, y, w, h, src: Path | str):
    """Inserisce un'immagine adattata a w×h, mantenendo aspect ratio
    se l'immagine non corrisponde."""
    return slide.shapes.add_picture(str(src), pt(x), pt(y),
                                    width=pt(w), height=pt(h))


def add_line(slide, x1, y1, x2, y2, color="Gold", weight_pt=6):
    """Linea retta colorata via shape FREEFORM è complesso: usiamo un
    Rectangle sottile come surrogate."""
    add_rect(slide, x1, y1 - weight_pt / 2, x2 - x1, weight_pt,
             fill_name=color)


# ────────── lucide icons → temp PNG ──────────


_ICON_CACHE: dict[tuple[str, int, str], Path] = {}
_TMP_DIR = ROOT / ".keynote_icon_cache"
_TMP_DIR.mkdir(exist_ok=True)


def _icon_png(name: str, size_px: int, color: str) -> Path:
    key = (name, size_px, color)
    if key in _ICON_CACHE:
        return _ICON_CACHE[key]
    rgb = COLORS[color]
    icon_img = render_icon(name, size=size_px * 4,
                           color=(rgb[0], rgb[1], rgb[2]),
                           stroke_width=size_px / 3)
    icon_img = icon_img.resize((size_px, size_px), Image.LANCZOS)
    out_path = _TMP_DIR / f"{name}_{size_px}_{color}.png"
    icon_img.save(out_path, "PNG", optimize=True)
    _ICON_CACHE[key] = out_path
    return out_path


def add_lucide(slide, name: str, cx: float, cy: float, size_pt: float,
               color="Ink"):
    """Inserisce un'icona lucide come PNG ad alta risoluzione."""
    px = int(round(size_pt * 4))  # 4× per nitidezza
    src = _icon_png(name, px, color)
    return add_image(slide, cx - size_pt / 2, cy - size_pt / 2,
                     size_pt, size_pt, src)


# ────────── header / footer ──────────


def header_bar(slide, page, total, title, kicker=None):
    add_image(slide, 50, 65, 60, 60, ICON_PATH)
    if kicker:
        add_text(slide, 140, 65, 700, 26, kicker.upper(),
                 size=18, color="Gold", bold=True)
    add_text(slide, 140, 95, 1300, 60, title,
             size=44, color="Ink", bold=True)
    # gold underline
    add_rect(slide, 140, 165, 80, 5, fill_name="Gold")
    add_text(slide, SLIDE_W_PT - 250, 95, 200, 26,
             f"{page:02d} / {total:02d}",
             size=16, color="Neutral", align="right")


def footer_bar(slide, page, total):
    # bottom border
    add_rect(slide, 0, SLIDE_H_PT - 60, SLIDE_W_PT, 1,
             fill_name="PaperDark")
    add_image(slide, 50, SLIDE_H_PT - 50, 30, 30, ICON_PATH)
    add_text(slide, 95, SLIDE_H_PT - 48, 800, 26,
             "Cadenza · Per i direttori dei conservatori",
             size=14, color="Ink", bold=True)
    add_text(slide, SLIDE_W_PT - 250, SLIDE_H_PT - 48, 200, 26,
             f"{page:02d} / {total:02d}",
             size=14, color="Neutral", align="right")


def page_bg(slide, color="Paper"):
    add_rect(slide, 0, 0, SLIDE_W_PT, SLIDE_H_PT, fill_name=color)


# ────────── slides ──────────


def slide_01_cover(slide, idx, total):
    page_bg(slide, "NavyDeep")
    add_rect(slide, 0, 0, SLIDE_W_PT, SLIDE_H_PT, fill_name="Navy")
    # logo center top
    icon_size = 260
    add_image(slide, (SLIDE_W_PT - icon_size) / 2, SLIDE_H_PT * 0.18,
              icon_size, icon_size, ICON_PATH)
    add_text(slide, 0, SLIDE_H_PT * 0.48, SLIDE_W_PT, 200,
             "CADENZA", size=140, color="White", bold=True, align="center")
    add_line(slide, SLIDE_W_PT / 2 - 130, SLIDE_H_PT * 0.62,
             SLIDE_W_PT / 2 + 130, SLIDE_H_PT * 0.62, color="Gold", weight_pt=6)
    add_text(slide, 0, SLIDE_H_PT * 0.66, SLIDE_W_PT, 70,
             "Software gratuito per il Conservatorio",
             size=42, color="GoldLight", bold=True, align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.74, SLIDE_W_PT, 50,
             "Booking aule · Monte Ore · Strumenti · Avvisi · Kiosk",
             size=24, color="Dim", align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.84, SLIDE_W_PT, 50,
             "Sviluppato da Danilo Russo, docente del Conservatorio",
             size=22, color="GoldPale", bold=True, align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.91, SLIDE_W_PT, 40,
             "30 aprile 2026 · Presentazione per Direzione, DSGA e responsabili IT",
             size=18, color="Dim", align="center")


def slide_02_perche(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Perché parliamo di questo, oggi", kicker="INTRO")
    cards = [
        ("21 mld €", "PNRR Missione 4",
         "Stanziati per la digitalizzazione delle istituzioni AFAM entro il 2026."),
        ("79", "Conservatori statali",
         "Più 50 istituti AFAM. Pubblica amministrazione con vincoli GDPR, MEPA, ANIS, conservazione sostitutiva."),
        ("3-15k€", "Budget software gestionale",
         "ASIMUT costa 15.000-40.000 €/anno. Le alternative italiane non sono verticali sul Conservatorio."),
    ]
    cw, gap = 540, 30
    base_x = (SLIDE_W_PT - 3 * cw - 2 * gap) / 2
    by = 250
    for i, (big, label, body) in enumerate(cards):
        x = base_x + i * (cw + gap)
        add_rect(slide, x, by, cw, 600, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=22)
        add_text(slide, x, by + 90, cw, 150, big,
                 size=100, color="Gold", bold=True, align="center")
        add_line(slide, x + cw / 2 - 60, by + 250, x + cw / 2 + 60, by + 250,
                 color="Gold", weight_pt=4)
        add_text(slide, x, by + 280, cw, 50, label,
                 size=28, color="Ink", bold=True, align="center")
        add_text(slide, x + 30, by + 360, cw - 60, 220, body,
                 size=20, color="Ink", align="center", line_spacing=1.4)
    footer_bar(slide, idx, total)


def slide_03_uno_pagina(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Cadenza in una pagina", kicker="COSA È")
    add_text(slide, 80, 230, SLIDE_W_PT - 160, 60,
             "Una piattaforma open-source progettata",
             size=42, color="Ink", bold=True)
    add_text(slide, 80, 290, SLIDE_W_PT - 160, 60,
             "specificamente per i conservatori italiani.",
             size=32, color="Navy")
    metrics = [
        ("31",   "modelli Sequelize",     "Architettura solida verificata"),
        ("100+", "endpoint API",          "REST coerente, documentazione fluente"),
        ("169",  "test automatici",       "Continuous integration GitHub Actions"),
        ("3",    "lingue supportate",     "Italiano · Inglese · Spagnolo"),
        ("1",    "deploy command",        "VPS Ubuntu o Docker, zero lock-in cloud"),
        ("0",    "doppie prenotazioni",   "Garantito a livello DB Postgres EXCLUDE"),
    ]
    by = 380
    cw = (SLIDE_W_PT - 200 - 30 * 2) / 3
    for i, (val, lbl, sub) in enumerate(metrics):
        col, row = i % 3, i // 3
        x = 90 + col * (cw + 30)
        y = by + row * 240
        add_rect(slide, x, y, cw, 210, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=18)
        add_text(slide, x + 30, y + 30, cw - 60, 100, val,
                 size=64, color="Navy", bold=True)
        add_text(slide, x + 30, y + 130, cw - 60, 40, lbl,
                 size=22, color="Ink", bold=True)
        add_text(slide, x + 30, y + 175, cw - 60, 30, sub,
                 size=16, color="Neutral")
    footer_bar(slide, idx, total)


def slide_04_4_risposte(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Quattro risposte concrete",
               kicker="DOMINI FUNZIONALI")
    items = [
        ("calendar-plus", "Booking aule",
         "Self-service in 3 tap. Anti-overlap garantito DB.\n"
         "Approval workflow per sale concerti. Waitlist auto-promote.",
         "•  /booking · /rooms · /my-bookings", "NavyLight"),
        ("clock", "Monte Ore docenti",
         "Workflow contrattuale del Conservatorio: vincoli 2-4 giorni\n"
         "/ settimana, soglia 324h/anno, sospensioni didattiche.",
         "•  /monte-ore · pannello docente + admin", "Gold"),
        ("package", "Inventario strumenti",
         "Catalogo + prestiti (5 stati). Reminder T-2gg, auto-overdue.\n"
         "PDF di consegna, email transazionali.",
         "•  /instruments · /my-loans", "Green"),
        ("megaphone", "Avvisi & Kiosk",
         "Bacheca con audience filter (ruolo · corso · edificio).\n"
         "Display di sala con rotazione concerti + annunci.",
         "•  /announcements · /display", "Violet"),
    ]
    by, cw, rh = 240, (SLIDE_W_PT - 200) / 2, 350
    for i, (icon_name, title, body, route, color) in enumerate(items):
        col, row = i % 2, i // 2
        x = 90 + col * (cw + 20)
        y = by + row * (rh + 20)
        add_rect(slide, x, y, cw, rh, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=18)
        add_rect(slide, x, y, 8, rh, fill_name=color, radius_pt=4)
        # icon badge top-right
        bx, by_b = x + cw - 100, y + 30
        add_rect(slide, bx, by_b, 70, 70, fill_name=color, radius_pt=14)
        add_lucide(slide, icon_name, bx + 35, by_b + 35, 44, color="White")
        add_text(slide, x + 40, y + 40, cw - 180, 50, title,
                 size=32, color="Ink", bold=True)
        add_text(slide, x + 40, y + 130, cw - 70, 100, body,
                 size=20, color="Ink", line_spacing=1.4)
        add_text(slide, x + 40, y + rh - 60, cw - 70, 30, route,
                 size=18, color="Navy", bold=True)
    footer_bar(slide, idx, total)


def slide_screenshot(slide, idx, total, kicker, title, png, caption):
    page_bg(slide)
    header_bar(slide, idx, total, title, kicker=kicker)
    img_x, img_y = 90, 215
    img_w, img_h = SLIDE_W_PT - 180, 760
    # frame
    add_rect(slide, img_x, img_y, img_w, img_h, fill_name="White",
             stroke_name="PaperDark", stroke_w_pt=2, radius_pt=14)
    # screenshot adattata mantenendo aspect: la PNG webapp è 1920×1080
    # quindi entra perfettamente; se mai non lo fosse, python-pptx
    # ridimensiona allo spazio dato.
    add_image(slide, img_x + 12, img_y + 12, img_w - 24, img_h - 24,
              SHOTS / png)
    add_text(slide, 80, SLIDE_H_PT - 95, SLIDE_W_PT - 160, 30, caption,
             size=18, color="Neutral", align="center")
    footer_bar(slide, idx, total)


def slide_05(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /dashboard", "Dashboard docente", "dashboard.png",
        "Quattro KPI personali, calendario aule giornaliero con drag-to-create, agenda prossime prenotazioni.")

def slide_06(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /rooms", "Aule del Conservatorio", "rooms.png",
        "Catalogo navigabile per edificio, capienza, dotazione. Filtri per attrezzatura, tipologia, prenotabilità.")

def slide_07(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /booking", "Prenota un'aula", "booking.png",
        "Timeline 30' (08-22), legenda tipologia (studio · lezione · prova · concerto), drag-to-select.")

def slide_08(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /admin/monte-ore", "Monte Ore docenti — pannello admin",
        "admin-monte-ore.png",
        "Workflow proposte annuali → approvazione coordinatore → generazione prenotazioni ricorrenti.")

def slide_09(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /admin/analytics", "Analytics direzione",
        "admin-analytics.png",
        "Heatmap occupazione 7×24, top aule, no-show rate, trend 8 settimane, export CSV/PDF.")

def slide_10(slide, idx, total):
    slide_screenshot(slide, idx, total,
        "SCHERMATA REALE · /admin/structure", "Struttura del Conservatorio",
        "admin-structure.png",
        "Anagrafica istituti, edifici, aule e catalogo dotazioni — tutto in un'unica pagina.")


def make_table(slide, headers, rows, x, y, total_w, *,
               header_fill="NavyDeep", row_h=64,
               col_widths=None, cell_styles=None):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [total_w / n_cols] * n_cols
    # header row
    add_rect(slide, x, y, total_w, row_h, fill_name=header_fill, radius_pt=10)
    cx = x
    for j, h in enumerate(headers):
        add_text(slide, cx, y, col_widths[j], row_h, h,
                 size=18, color="White", bold=True, align="center")
        cx += col_widths[j]
    for i, row in enumerate(rows):
        ry = y + (i + 1) * row_h
        if i % 2 == 0:
            add_rect(slide, x, ry, total_w, row_h, fill_name="White")
        cx = x
        for j, cell in enumerate(row):
            color = "Ink"
            bold = False
            if cell == "—":
                color = "Neutral"
            elif cell_styles and cell_styles[i][j]:
                style = cell_styles[i][j]
                if style == "pos":
                    color = "GreenDark"; bold = True
                elif style == "red":
                    color = "RedDark"; bold = True
            align = "left" if j == 0 else "center"
            padding = 16 if j == 0 else 0
            add_text(slide, cx + padding, ry, col_widths[j] - 2 * padding, row_h,
                     cell, size=18, color=color, bold=bold, align=align)
            cx += col_widths[j]


def slide_11_vs_asimut(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Cadenza vs ASIMUT", kicker="CONFRONTO 1/2")
    headers = ["Aspetto", "ASIMUT", "Cadenza"]
    rows = [
        ["Pricing annuo (medio)",   "15.000-40.000 €",       "463-2.003 €"],
        ["Cloud / Self-host",       "Solo cloud (Danimarca)", "Cloud IT o self-host"],
        ["Lingua / supporto",       "Inglese, fuso CET",     "Italiano, made in Italy"],
        ["Monte Ore docenti AFAM",  "—",                     "✓ verticale italiana"],
        ["SPID / CIE",              "—",                     "Roadmap Sprint E"],
        ["ANIS / MIUR export",      "—",                     "Roadmap Sprint E"],
        ["Bot Telegram / WhatsApp", "—",                     "✓ production-ready"],
        ["Inventario strumenti",    "Generico",              "✓ workflow prestito"],
    ]
    cell_styles = [[None, None, "pos"] for _ in rows]
    make_table(slide, headers, rows, x=80, y=240, total_w=SLIDE_W_PT - 160,
               col_widths=[800, 510, 450], cell_styles=cell_styles)
    footer_bar(slide, idx, total)


def slide_12_vs_easystaff(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total,
               "Cadenza vs EasyStaff / EasyAcademy", kicker="CONFRONTO 2/2")
    headers = ["Aspetto", "EasyStaff", "Cadenza"]
    rows = [
        ["Verticalità Conservatori", "Pacchetto università generaliste",
         "Verticale AFAM dal giorno 1"],
        ["Monte Ore",                "—", "✓ workflow contrattuale"],
        ["Eventi / sale concerti",   "—", "✓ approval + iCal + kiosk"],
        ["Inventario strumenti",     "—", "✓ catalogo + prestiti completo"],
        ["Pricing AFAM medio",       "8.000-15.000 €", "463-2.003 €"],
        ["Onboarding tipico",        "8-12 settimane", "0,5-2 giornate"],
        ["Open-source / self-host",  "—", "✓ codice sorgente in licenza"],
        ["Lingua e supporto",        "Italiano OK", "Italiano + IT/EN/ES"],
    ]
    cell_styles = [[None, None, "pos"] for _ in rows]
    make_table(slide, headers, rows, x=80, y=240, total_w=SLIDE_W_PT - 160,
               col_widths=[800, 510, 450], cell_styles=cell_styles)
    footer_bar(slide, idx, total)


def slide_13_compliance(slide, idx, total):
    page_bg(slide, "NavyDeep")
    add_rect(slide, 0, 0, SLIDE_W_PT, SLIDE_H_PT, fill_name="Navy")
    add_text(slide, 0, 130, SLIDE_W_PT, 100, "Italiano per definizione",
             size=64, color="White", bold=True, align="center")
    add_line(slide, SLIDE_W_PT / 2 - 80, 250, SLIDE_W_PT / 2 + 80, 250,
             color="Gold", weight_pt=5)
    add_text(slide, 0, 290, SLIDE_W_PT, 50,
             "Quattro garanzie che i competitor esteri non offrono nativamente.",
             size=24, color="Dim", align="center")
    badges = [
        ("flag", "Codice italiano", "Sviluppato in Italia, manutenzione in italiano",
         "Niente time-zone CET-only, niente lingua di supporto inglese"),
        ("shield-check", "GDPR Garante 06/2021", "Cookie, DPIA, ROPA, audit log",
         "Provvedimento 06/2021 e art. 13-17 GDPR coperti by-default"),
        ("scale", "Roadmap PA italiana", "SPID/CIE · PEC · ANIS/MIUR",
         "Sprint E pianificato, sviluppo on-demand al primo Conservatorio"),
        ("building-2", "MEPA-ready", "Pubblica amministrazione, MEPA, fattura elettronica",
         "Tutti i piani sotto soglia 75.000 € (affidamento diretto)"),
    ]
    bx, by = 200, 440
    bw = (SLIDE_W_PT - 400 - 60) / 2
    bh = 230
    for i, (icon_name, t1, t2, foot) in enumerate(badges):
        col, row = i % 2, i // 2
        x = bx + col * (bw + 60)
        y = by + row * (bh + 30)
        add_rect(slide, x, y, bw, bh, fill_name="NavyLight",
                 stroke_name="Gold", stroke_w_pt=2, radius_pt=20)
        # gold badge with icon
        add_rect(slide, x + 30, y + 60, 100, 100, fill_name="Gold", radius_pt=20)
        add_lucide(slide, icon_name, x + 80, y + 110, 64, color="NavyDeep")
        add_text(slide, x + 160, y + 50, bw - 180, 50, t1,
                 size=28, color="White", bold=True)
        add_text(slide, x + 160, y + 105, bw - 180, 30, t2,
                 size=18, color="Dim")
        add_text(slide, x + 160, y + 145, bw - 180, 60, foot,
                 size=16, color="GoldLight", line_spacing=1.4)


def slide_14_costi(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Costi reali per il Conservatorio", kicker="TUTTO QUI")
    add_text(slide, 80, 230, SLIDE_W_PT - 160, 60,
             "Software gratuito. Pagate solo l'infrastruttura.",
             size=36, color="Ink", bold=True)
    add_text(slide, 80, 295, SLIDE_W_PT - 160, 50,
             "Niente licenze, niente canoni, niente lock-in.",
             size=28, color="Navy")
    rows = [
        ("server", "VPS Hetzner CPX31",
         "4 vCPU · 8 GB · 160 GB SSD · datacenter EU", "192 €", "/anno"),
        ("globe", "Dominio + Let's Encrypt",
         "Sotto-dominio del Conservatorio. TLS gratuito.", "15 €", "/anno"),
        ("database", "Backup off-site",
         "Hetzner Storage Box 1 TB, 30 giorni retention", "36 €", "/anno"),
        ("sparkles", "Claude Pro (manutenzione AI-assisted)",
         "Per evoluzioni e bug fix. Opzionale: Max 1.760 €/anno", "220 €", "/anno"),
    ]
    by, rh = 380, 110
    for i, (icon_name, lbl, sub, price, unit) in enumerate(rows):
        y = by + i * (rh + 12)
        add_rect(slide, 80, y, SLIDE_W_PT - 160, rh, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=14)
        add_rect(slide, 80, y, 8, rh, fill_name="Gold")
        add_rect(slide, 110, y + 25, 60, 60, fill_name="GoldPale", radius_pt=12)
        add_lucide(slide, icon_name, 140, y + 55, 36, color="Gold")
        add_text(slide, 190, y + 18, SLIDE_W_PT - 560, 40, lbl,
                 size=26, color="Ink", bold=True)
        add_text(slide, 190, y + 65, SLIDE_W_PT - 560, 30, sub,
                 size=18, color="Neutral")
        add_text(slide, SLIDE_W_PT - 280, y + 22, 200, 60, price,
                 size=42, color="Navy", bold=True)
        add_text(slide, SLIDE_W_PT - 280, y + 80, 200, 30, unit,
                 size=18, color="Neutral")
    # total bar
    ty = by + 4 * (rh + 12) + 10
    add_rect(slide, 80, ty, SLIDE_W_PT - 160, 110, fill_name="NavyDeep", radius_pt=18)
    add_text(slide, 110, ty + 22, SLIDE_W_PT - 500, 40,
             "Totale annuo Conservatorio", size=28, color="White", bold=True)
    add_text(slide, 110, ty + 70, SLIDE_W_PT - 500, 30,
             "Tutto incluso, IVA esclusa", size=18, color="Dim")
    add_text(slide, SLIDE_W_PT - 320, ty + 26, 240, 70, "≈ 463 €",
             size=46, color="GoldLight", bold=True)
    footer_bar(slide, idx, total)


def slide_15_tco(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Quanto risparmiate davvero",
               kicker="TCO 1 / 5 / 10 ANNI")
    headers = ["Periodo", "ASIMUT (medio)", "Cadenza Pro", "Risparmio"]
    rows = [
        ["Anno 1",  "22.500 €",  "463 €",   "− 22.037 €"],
        ["3 anni",  "67.500 €",  "1.389 €", "− 66.111 €"],
        ["5 anni",  "112.500 €", "2.315 €", "− 110.185 €"],
        ["10 anni", "225.000 €", "4.630 €", "− 220.370 €"],
    ]
    cell_styles = [[None, "red", "pos", None] for _ in rows]
    make_table(slide, headers, rows, x=80, y=240,
               total_w=SLIDE_W_PT - 160,
               col_widths=[520, 400, 400, SLIDE_W_PT - 160 - 1320],
               row_h=70, cell_styles=cell_styles)
    # gold bubbles for risparmio
    bubble_x = 80 + 520 + 400 + 400 + 30
    bubble_w = (SLIDE_W_PT - 160 - 1320) - 60
    for i, row in enumerate(rows):
        y = 240 + (i + 1) * 70 + 8
        add_rect(slide, bubble_x, y, bubble_w, 70 - 16,
                 fill_name="Gold", radius_pt=14)
        add_text(slide, bubble_x, y, bubble_w, 70 - 16, row[3],
                 size=22, color="NavyDeep", bold=True, align="center")
    # callout
    cy = 240 + 5 * 70 + 60
    add_rect(slide, 80, cy, SLIDE_W_PT - 160, 280, fill_name="White",
             stroke_name="Gold", stroke_w_pt=3, radius_pt=18)
    add_text(slide, 110, cy + 24, SLIDE_W_PT - 220, 50,
             "Su 10 anni, 220.000 € equivalgono a:", size=26, color="Ink", bold=True)
    bullets = [
        "•  ~ 4 stipendi annuali docente di II fascia (lordo amministrazione)",
        "•  ~ 2 organi a canne medi · oppure 6-8 pianoforti gran coda",
        "•  ~ 40-50 borse di studio annuali per studenti meritevoli (5.000 € cad.)",
        "•  ~ 10 ristrutturazioni di un'aula completa di insonorizzazione e impianto audio",
    ]
    for j, b in enumerate(bullets):
        add_text(slide, 120, cy + 80 + j * 46, SLIDE_W_PT - 240, 40, b,
                 size=20, color="Ink")
    footer_bar(slide, idx, total)


def slide_16_vps(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Dove ospitare Cadenza",
               kicker="VPS · FINO A 5.000 UTENTI")
    headers = ["Provider", "Datacenter",
               "Piano (1.5-3k ut.)", "Piano (3-5k ut.)", "Note"]
    rows = [
        ["Hetzner Cloud",     "DE / FI", "CPX31 — 16 €/mese", "CPX41 — 28 €/mese",
         "★ Miglior prezzo/prest."],
        ["Hetzner Dedicated", "DE",       "CCX23 — 39 €/mese", "CCX33 — 78 €/mese",
         "Carichi predicibili"],
        ["OVHcloud",          "FR / DE",  "VPS Comfort — 15 €/mese",
         "VPS Elite — 29 €/mese",          "Provider EU consolidato"],
        ["Ionos Cloud",       "DE / IT",  "Cloud L — 25 €/mese",
         "Cloud XL — 50 €/mese",           "Sovranità EU, supp. IT"],
        ["Aruba Cloud",       "Italia",   "Smart VS3 — 22 €/mese",
         "Smart VS4 — 50 €/mese",          "★ Sovranità italiana"],
        ["Register.it",       "Italia",   "Cloud M — 30 €/mese",
         "Cloud L — 60 €/mese",            "MEPA-ready, fattura PA"],
        ["DigitalOcean",      "DE / NL",  "Basic 4/8 — 48 €/mese",
         "Premium 4/8 — 56 €/mese",        "Ecosystem mature"],
    ]
    make_table(slide, headers, rows, x=80, y=230,
               total_w=SLIDE_W_PT - 160,
               col_widths=[230, 200, 380, 380, SLIDE_W_PT - 160 - 1190],
               row_h=60)
    cy = 230 + 8 * 60 + 30
    add_rect(slide, 80, cy, SLIDE_W_PT - 160, 110,
             fill_name="NavyDeep", radius_pt=18)
    add_text(slide, 80, cy + 25, SLIDE_W_PT - 160, 40,
             "Tutti gli scenari rientrano in affidamento diretto (D.Lgs. 36/2023, art. 50)",
             size=22, color="White", bold=True, align="center")
    add_text(slide, 80, cy + 70, SLIDE_W_PT - 160, 30,
             "Spesa annua totale: 451 € (OVH) — 738 € (Hetzner dedicated) · "
             "Tutto in fattura PA elettronica.",
             size=18, color="GoldLight", align="center")
    footer_bar(slide, idx, total)


def slide_17_attivazione(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Cosa serve per partire", kicker="ATTIVAZIONE")
    add_text(slide, 80, 230, SLIDE_W_PT - 160, 50,
             "Decisione → operatività in 2-4 settimane.",
             size=32, color="Ink", bold=True)
    items = [
        ("server", "VPS + dominio",
         "Acquisto del server (Hetzner / Aruba / Ionos / OVH).\nA carico del Conservatorio. ~ 200 €/anno."),
        ("terminal", "Provisioning",
         "Install pacchetti, deploy Cadenza, certificato TLS.\nA carico dell'autore. 2-3 ore una sola volta."),
        ("upload", "Import anagrafica",
         "Utenti, aule, edifici, dotazioni da Excel/CSV.\nAutore + DSGA. 1-2 ore."),
        ("palette", "Branding",
         "Logo del Conservatorio, colori, intestazioni.\nAutore. 30 minuti."),
        ("users", "Formazione admin",
         "Sessione 90 minuti remoto per Direzione/DSGA/IT.\nAutore. Registrabile."),
        ("play-circle", "Formazione utenti",
         "Sessione 60 minuti per docenti / studenti pilota.\nAutore. Materiale video registrato."),
    ]
    by, cw, rh = 320, (SLIDE_W_PT - 200) / 2, 150
    for i, (icon_name, kicker, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = 90 + col * (cw + 20)
        y = by + row * (rh + 16)
        add_rect(slide, x, y, cw, rh, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=14)
        add_rect(slide, x, y, 6, rh, fill_name="Gold")
        add_rect(slide, x + 20, y + 30, 80, 80, fill_name="NavyDeep", radius_pt=14)
        add_lucide(slide, icon_name, x + 60, y + 70, 48, color="GoldLight")
        add_text(slide, x + 120, y + 22, cw - 140, 40, kicker,
                 size=24, color="Navy", bold=True)
        add_text(slide, x + 120, y + 70, cw - 140, 80, body,
                 size=18, color="Ink", line_spacing=1.4)
    add_rect(slide, 80, 920, SLIDE_W_PT - 160, 70, fill_name="NavyDeep", radius_pt=18)
    add_text(slide, 80, 940, SLIDE_W_PT - 160, 40,
             "Nessun contratto di licenza. Codice sorgente del Conservatorio.",
             size=22, color="GoldLight", bold=True, align="center")
    footer_bar(slide, idx, total)


def slide_18_roadmap(slide, idx, total):
    page_bg(slide)
    header_bar(slide, idx, total, "Roadmap 12 mesi", kicker="PROSSIMI PASSI")
    sprints = [
        ("sparkles",       "Sprint A", "UX quick wins",
         "Push notif · iframe concerti · Privacy display · Card avvisi", "Q2 2026"),
        ("bar-chart-3",    "Sprint B", "Analytics +",
         "Maintenance schedule · Report email YoY analytics", "Q2 2026"),
        ("calendar-check", "Sprint C", "Task mgmt eventi",
         "Gap parità ASIMUT (staff workflow)", "Q3 2026"),
        ("wrench",         "Sprint D", "Tech debt",
         "Docker compose · Sequelize-CLI · Coverage 70%", "Q3 2026"),
        ("shield-check",   "Sprint E", "PA italiana",
         "SPID/CIE · PEC · ANIS/MIUR · Conservazione", "Q4 2026"),
        ("bot",            "Sprint F", "Bot completi",
         "WhatsApp Cloud · Signal · Email IMAP poller", "Q1 2027"),
    ]
    by = 250
    cw = (SLIDE_W_PT - 180 - 30 * 5) / 6
    for i, (icon_name, k, t1, body, when) in enumerate(sprints):
        x = 90 + i * (cw + 30)
        add_rect(slide, x, by, cw, 540, fill_name="White",
                 stroke_name="PaperDark", stroke_w_pt=2, radius_pt=16)
        add_rect(slide, x, by, cw, 70, fill_name="NavyDeep", radius_pt=16)
        add_text(slide, x, by + 18, cw, 40, k,
                 size=20, color="GoldLight", bold=True, align="center")
        add_oval(slide, x + cw / 2 - 40, by + 100, 80, 80, fill_name="Gold")
        add_lucide(slide, icon_name, x + cw / 2, by + 140, 50, color="NavyDeep")
        add_text(slide, x, by + 220, cw, 40, t1,
                 size=22, color="Ink", bold=True, align="center")
        add_text(slide, x + 10, by + 290, cw - 20, 160, body,
                 size=16, color="Ink", align="center", line_spacing=1.4)
        add_text(slide, x, by + 480, cw, 40, when,
                 size=18, color="Gold", bold=True, align="center")
    footer_bar(slide, idx, total)


def slide_19_cta(slide, idx, total):
    page_bg(slide, "NavyDeep")
    add_rect(slide, 0, 0, SLIDE_W_PT, SLIDE_H_PT, fill_name="Navy")
    icon_size = 220
    add_image(slide, (SLIDE_W_PT - icon_size) / 2, SLIDE_H_PT * 0.18,
              icon_size, icon_size, ICON_PATH)
    add_text(slide, 0, SLIDE_H_PT * 0.46, SLIDE_W_PT, 200, "CADENZA",
             size=130, color="White", bold=True, align="center")
    add_line(slide, SLIDE_W_PT / 2 - 130, SLIDE_H_PT * 0.59,
             SLIDE_W_PT / 2 + 130, SLIDE_H_PT * 0.59, color="Gold", weight_pt=6)
    add_text(slide, 0, SLIDE_H_PT * 0.65, SLIDE_W_PT, 70,
             "Un dono al Conservatorio.",
             size=42, color="GoldLight", bold=True, align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.72, SLIDE_W_PT, 70,
             "Per i prossimi dieci anni almeno.",
             size=42, color="Dim", bold=True, align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.83, SLIDE_W_PT, 50,
             "Danilo Russo · docente del Conservatorio",
             size=24, color="GoldPale", bold=True, align="center")
    add_text(slide, 0, SLIDE_H_PT * 0.91, SLIDE_W_PT, 40,
             "Una demo dal vivo, in italiano: 30 minuti per mostrarvi tutto.",
             size=20, color="Dim", align="center")


SLIDE_FUNCS = [
    slide_01_cover, slide_02_perche, slide_03_uno_pagina, slide_04_4_risposte,
    slide_05, slide_06, slide_07, slide_08, slide_09, slide_10,
    slide_11_vs_asimut, slide_12_vs_easystaff, slide_13_compliance,
    slide_14_costi, slide_15_tco, slide_16_vps, slide_17_attivazione,
    slide_18_roadmap, slide_19_cta,
]


def main():
    prs = Presentation()
    # 16:9 widescreen 1920×1080 pt
    prs.slide_width = pt(SLIDE_W_PT)
    prs.slide_height = pt(SLIDE_H_PT)

    blank_layout = prs.slide_layouts[6]   # "Blank"
    n = len(SLIDE_FUNCS)
    print(f"Building Keynote (.pptx) with {n} slides…")
    for i, fn in enumerate(SLIDE_FUNCS):
        slide = prs.slides.add_slide(blank_layout)
        # rimuove eventuali placeholder rimasti (titolo, ecc.)
        for shp in list(slide.shapes):
            if shp.has_text_frame and shp.placeholder_format is not None:
                sp = shp._element
                sp.getparent().remove(sp)
        fn(slide, i + 1, n)
        print(f"  ✓ slide {i + 1:02d}")
    prs.save(str(OUT))
    print(f"\n.pptx pronto: {OUT} ({OUT.stat().st_size / 1024 / 1024:.1f} MB)")
    print(
        "\nPer convertirlo in formato Keynote (.key):\n"
        "  1. Apri il .pptx con Keynote (doppio click su Mac con Keynote installato)\n"
        "  2. File → Salva come… → formato Keynote (.key)\n"
        "Keynote importa automaticamente shape, testo, immagini, font."
    )


if __name__ == "__main__":
    main()
